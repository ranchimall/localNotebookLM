"""
Turn RAG query results into a PowerPoint or PDF -- one slide/page per
retrieved note, raw text, no summarization.

Usage:
    python present.py pptx "your question" --top-k 50 --output my_deck
    python present.py pdf  "your question" --top-k 50 --output my_deck
"""

import argparse

import chromadb
from chromadb.utils import embedding_functions
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.units import inch
from reportlab.lib.utils import simpleSplit
from reportlab.pdfgen import canvas

from query import retrieve  # reuse the existing retrieval logic

PDF_PAGE_SIZE = landscape(letter)
PDF_MARGIN = 0.6 * inch
PDF_LINE_HEIGHT = 16


# ---------- shared: retrieval ----------

def get_results(args):
    embed_fn = embedding_functions.SentenceTransformerEmbeddingFunction(model_name=args.model)
    client = chromadb.PersistentClient(path=args.db_path)
    collection = client.get_collection(name=args.collection, embedding_function=embed_fn)
    results = retrieve(collection, args.query, n_results=args.top_k, label_filter=args.label)
    n_found = len(results["documents"][0])
    print(f"Retrieved {n_found} notes for: {args.query}")
    return results


# ---------- PPTX building ----------

def build_pptx(results, query_text, show_full=False):
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = query_text
    title_slide.placeholders[1].text = f"{len(results['documents'][0])} retrieved notes"

    docs = results["documents"][0]
    metas = results["metadatas"][0]
    dists = results["distances"][0]
    layout = prs.slide_layouts[1]  # "Title and Content"

    for doc, meta, dist in zip(docs, metas, dists):
        similarity = 1 - dist
        title = meta.get("title", "(untitled)")
        source = meta.get("source_file", "?")
        text = meta.get("full_text", doc) if show_full else doc

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        body = slide.placeholders[1].text_frame
        body.word_wrap = True
        body.text = text
        font_size = 14 if len(text) < 500 else 11
        for paragraph in body.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)

        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        footer.text_frame.text = f"similarity={similarity:.3f} | source={source}"
        footer.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        footer.text_frame.paragraphs[0].runs[0].font.italic = True

    return prs


# ---------- PDF building ----------

def _pdf_draw_footer(c, width, similarity, source):
    c.setFont("Helvetica-Oblique", 9)
    c.drawString(PDF_MARGIN, PDF_MARGIN / 2, f"similarity={similarity:.3f} | source={source}")


def build_pdf(pdf_path, query_text, results, show_full=False):
    docs = results["documents"][0]
    metas = results["metadatas"][0]
    dists = results["distances"][0]
    width, height = PDF_PAGE_SIZE

    c = canvas.Canvas(pdf_path, pagesize=PDF_PAGE_SIZE)

    c.setFont("Helvetica-Bold", 24)
    c.drawCentredString(width / 2, height / 2 + 20, query_text)
    c.setFont("Helvetica", 14)
    c.drawCentredString(width / 2, height / 2 - 10, f"{len(docs)} retrieved notes")
    c.showPage()

    for doc, meta, dist in zip(docs, metas, dists):
        similarity = 1 - dist
        title = meta.get("title", "(untitled)")
        source = meta.get("source_file", "?")
        text = meta.get("full_text", doc) if show_full else doc

        c.setFont("Helvetica-Bold", 18)
        c.drawString(PDF_MARGIN, height - PDF_MARGIN, title)

        c.setFont("Helvetica", 12)
        lines = []
        for paragraph in text.split("\n"):
            lines.extend(simpleSplit(paragraph, "Helvetica", 12, width - 2 * PDF_MARGIN) or [""])

        y = height - PDF_MARGIN - 40
        for line in lines:
            if y < PDF_MARGIN + 30:
                _pdf_draw_footer(c, width, similarity, source)
                c.showPage()
                c.setFont("Helvetica-Bold", 14)
                c.drawString(PDF_MARGIN, height - PDF_MARGIN, f"{title} (cont.)")
                c.setFont("Helvetica", 12)
                y = height - PDF_MARGIN - 40
            c.drawString(PDF_MARGIN, y, line)
            y -= PDF_LINE_HEIGHT

        _pdf_draw_footer(c, width, similarity, source)
        c.showPage()

    c.save()


# ---------- CLI ----------

def add_common_args(subparser):
    subparser.add_argument("query", help="Your question")
    subparser.add_argument("--db-path", default="./chroma_db")
    subparser.add_argument("--collection", default="keep_notes")
    subparser.add_argument("--model", default="all-MiniLM-L6-v2")
    subparser.add_argument("--top-k", type=int, default=5)
    subparser.add_argument("--label", default=None)
    subparser.add_argument("--full", action="store_true", help="Use full original note instead of just the matched chunk")
    subparser.add_argument("--output", default="rag_results", help="Output filename, no extension")


def main():
    parser = argparse.ArgumentParser(description="Turn RAG query results into a PPTX or PDF")
    subparsers = parser.add_subparsers(dest="command", required=True)

    pptx_parser = subparsers.add_parser("pptx", help="Generate a .pptx deck")
    add_common_args(pptx_parser)

    pdf_parser = subparsers.add_parser("pdf", help="Generate a .pdf")
    add_common_args(pdf_parser)

    args = parser.parse_args()
    results = get_results(args)

    if args.command == "pptx":
        prs = build_pptx(results, args.query, show_full=args.full)
        path = f"{args.output}.pptx"
        prs.save(path)
        print(f"Saved {path} ({len(prs.slides)} slides)")

    elif args.command == "pdf":
        path = f"{args.output}.pdf"
        build_pdf(path, args.query, results, show_full=args.full)
        print(f"Saved {path}")


if __name__ == "__main__":
    main()